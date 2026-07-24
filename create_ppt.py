from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 宽屏
prs.slide_height = Inches(7.5)

# 定义颜色方案
PRIMARY = RGBColor(0x1F, 0x3A, 0x5F)  # 深蓝色
SECONDARY = RGBColor(0xC8, 0xD3, 0xE6)  # 浅蓝色
ACCENT = RGBColor(0x0D, 0x94, 0x88)  # 青色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1F, 0x1F, 0x1F)
GRAY = RGBColor(0x7A, 0x7A, 0x7A)
LIGHT_BG = RGBColor(0xF7, 0xF5, 0xF0)  # 米白色背景

def add_background(slide, color):
    """为幻灯片添加背景色"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()  # 无边框
    # 将背景移到最底层
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    """添加文本框"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=BLACK):
    """添加项目符号列表"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return tb

def add_section_header(slide, title):
    """添加章节标题"""
    add_text_box(slide, 0.5, 0.5, 12.3, 1.0, title, font_size=32, bold=True, color=PRIMARY)
    # 添加装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

# 1. 封面幻灯片
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
add_background(slide, PRIMARY)
add_text_box(slide, 1, 2, 11.3, 2, "英语学习小程序", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4, 11.3, 1, "EnglishMaster - 从零基础到自由沟通", font_size=24, color=SECONDARY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11.3, 1, "产品方案与技术架构", font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# 2. 产品概述
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_section_header(slide, "产品概述")

add_text_box(slide, 0.5, 1.8, 12.3, 1, "产品定位", font_size=24, bold=True, color=PRIMARY)
add_text_box(slide, 0.5, 2.5, 12.3, 1.5, "面向国内零基础成人，打造从音标发音矫正→单词积累→系统语法→短句情景对话→自由和AI外国人闲聊的沉浸式自学口语闭环。", font_size=18)

add_text_box(slide, 0.5, 4, 6, 1, "目标用户", font_size=20, bold=True, color=PRIMARY)
add_bullet_list(slide, 0.5, 4.5, 6, 2.5, [
    "英语零基础或基础薄弱的成人学习者",
    "希望提高口语发音和流利度的职场人士",
    "需要系统学习语法、应对日常沟通的学习者"
], font_size=16)

add_text_box(slide, 7, 4, 6, 1, "核心价值", font_size=20, bold=True, color=PRIMARY)
add_bullet_list(slide, 7, 4.5, 6, 2.5, [
    "AI纠音：音素级发音评估，纠正中式发音",
    "科学记忆：艾宾浩斯记忆曲线，高效背单词",
    "场景化语法：融入日常对话，学完就能用",
    "AI对话陪练：模拟真实场景，锻炼口语反应"
], font_size=16)

# 3. 技术平台选型
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_section_header(slide, "技术平台选型")

add_text_box(slide, 0.5, 1.8, 12.3, 1, "最终选型：微信小程序（核心主产品）", font_size=24, bold=True, color=PRIMARY)

add_text_box(slide, 0.5, 2.8, 6, 1, "技术栈", font_size=20, bold=True, color=PRIMARY)
add_bullet_list(slide, 0.5, 3.3, 6, 3, [
    "前端框架：UniApp/Taro（Vue语法，可跨平台编译）",
    "后端服务：Node.js/Java + 云服务器",
    "语音评测：微软Azure发音评测 / 讯飞语音口语评测SDK",
    "数据存储：MySQL + Redis + 对象存储（OSS）",
    "推送服务：微信小程序订阅消息"
], font_size=16)

add_text_box(slide, 7, 2.8, 6, 1, "选型理由", font_size=20, bold=True, color=PRIMARY)
add_bullet_list(slide, 7, 3.3, 6, 3, [
    "国内用户体量最大，打开零门槛",
    "微信原生录音API成熟，完美兼容语音评测SDK",
    "社交裂变无敌：打卡分享、好友组队背单词",
    "服务通知可定时推送背单词提醒",
    "微信支付成熟，后期开通会员简单",
    "前期低成本启动、快速迭代"
], font_size=16)

# 4. 核心功能模块
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_section_header(slide, "核心功能模块")

# 创建三个功能卡片
def add_feature_card(slide, x, y, title, items):
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(4.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = SECONDARY
    card.line.width = Pt(1)
    
    # 标题
    add_text_box(slide, x + 0.2, y + 0.2, 3.4, 0.5, title, font_size=20, bold=True, color=PRIMARY)
    
    # 内容
    add_bullet_list(slide, x + 0.2, y + 0.8, 3.4, 3.5, items, font_size=14)

add_feature_card(slide, 0.5, 1.8, "发音跟读模块", [
    "音标发音通关课",
    "48个英美音标逐个教学",
    "AI发音评估与纠错",
    "单词逐词跟读",
    "短句/段落影子跟读",
    "AI情景对话陪练"
])

add_feature_card(slide, 4.7, 1.8, "单词记忆模块", [
    "分级词库（500/1500/3000词）",
    "艾宾浩斯记忆曲线",
    "听音辨义、拼写默写",
    "图片联想记忆",
    "例句语境学习",
    "错题生词本"
])

add_feature_card(slide, 8.9, 1.8, "语法学习模块", [
    "三阶段语法课程体系",
    "场景化语法讲解",
    "填空练习、造句练习",
    "AI批改语法错误",
    "实时语法纠错",
    "语法知识图谱"
])

# 5. 发音跟读模块详解
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_section_header(slide, "发音跟读模块详解")

add_text_box(slide, 0.5, 1.8, 12.3, 1, "功能流程", font_size=24, bold=True, color=PRIMARY)

# 创建流程图
def add_flow_step(slide, x, y, step_num, title, description):
    # 步骤圆圈
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.8), Inches(0.8)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    
    # 步骤数字
    add_text_box(slide, x + 0.1, y + 0.15, 0.6, 0.5, str(step_num), font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # 标题和描述
    add_text_box(slide, x - 0.3, y + 1, 1.4, 0.5, title, font_size=14, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x - 0.5, y + 1.5, 1.8, 1, description, font_size=12, color=BLACK, alignment=PP_ALIGN.CENTER)

# 添加流程步骤
flow_steps = [
    (1, "选择课程", "音标/单词/句子"),
    (2, "播放标准发音", "可重复播放"),
    (3, "用户跟读录音", "实时录制"),
    (4, "AI语音评测", "音素级分析"),
    (5, "反馈与纠错", "评分+改进建议")
]

for i, (num, title, desc) in enumerate(flow_steps):
    x = 0.8 + i * 2.5
    add_flow_step(slide, x, 3, num, title, desc)
    # 连接箭头（用矩形模拟）
    if i < len(flow_steps) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x + 1.8), Inches(3.35), Inches(0.5), Inches(0.1)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY
        arrow.line.fill.background()

add_text_box(slide, 0.5, 5.5, 12.3, 1, "技术实现", font_size=20, bold=True, color=PRIMARY)
add_bullet_list(slide, 0.5, 6, 12.3, 1.5, [
    "语音录制：使用wx.getRecorderManager录制用户发音",
    "语音评测：集成微软Azure发音评测SDK，音素级评估",
    "反馈展示：显示评分、识别文本、标准文本，高亮错误音节",
    "多级支持：音标→单词→短句→段落→情景对话"
], font_size=16)

# 6. 单词记忆模块详解
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_section_header(slide, "单词记忆模块详解")

add_text_box(slide, 0.5, 1.8, 6, 1, "记忆闭环流程", font_size=24, bold=True, color=PRIMARY)

# 添加记忆闭环图
def add_memory_cycle(slide, x, y, step, title, color):
    # 圆形步骤
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(1.2), Inches(1.2)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    
    # 步骤文字
    add_text_box(slide, x + 0.1, y + 0.3, 1, 0.6, step, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # 标题
    add_text_box(slide, x - 0.2, y + 1.3, 1.6, 0.5, title, font_size=12, color=BLACK, alignment=PP_ALIGN.CENTER)

# 添加四个记忆步骤
colors = [PRIMARY, ACCENT, RGBColor(0xB8, 0x50, 0x42), RGBColor(0x2C, 0x5F, 0x2D)]
steps = [
    ("听音辨义", "播放发音选择释义"),
    ("看英文想中文", "显示单词回忆释义"),
    ("拼写默写", "根据释义拼写单词"),
    ("单词跟读", "跟读发音AI评分")
]

positions = [(1, 2.5), (4, 2.5), (1, 4.5), (4, 4.5)]
for i, (title, desc) in enumerate(steps):
    x, y = positions[i]
    add_memory_cycle(slide, x, y, f"第{i+1}步", title, colors[i])
    add_text_box(slide, x - 0.2, y + 2.6, 1.6, 0.5, desc, font_size=12, color=BLACK, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 7, 1.8, 6, 1, "SRS算法（间隔重复系统）", font_size=24, bold=True, color=PRIMARY)
add_bullet_list(slide, 7, 2.5, 6, 3, [
    "基于艾宾浩斯遗忘曲线",
    "根据用户记忆程度动态调整复习间隔",
    "SM-2算法实现",
    "记住的单词延长复习间隔",
    "忘记的单词缩短复习间隔",
    "智能安排每日复习计划"
], font_size=16)

add_text_box(slide, 7, 5, 6, 1, "辅助工具", font_size=20, bold=True, color=PRIMARY)
add_bullet_list(slide, 7, 5.5, 6, 2, [
    "错题生词本：自动收集标记单词",
    "薄弱词汇强化：针对错误率高的单词",
    "单词打卡日历：可视化学习连续性",
    "好友背单词PK：微信好友排名激励"
], font_size=14)

# 7. 语法学习模块详解
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_section_header(slide, "语法学习模块详解")

add_text_box(slide, 0.5, 1.8, 12.3, 1, "三阶段语法课程体系", font_size=24, bold=True, color=PRIMARY)

# 创建三个阶段卡片
def add_stage_card(slide, x, y, stage_num, title, items, color):
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(3)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    
    # 阶段标签
    add_text_box(slide, x + 0.2, y + 0.2, 1, 0.5, f"阶段{stage_num}", font_size=16, bold=True, color=WHITE)
    
    # 标题
    add_text_box(slide, x + 0.2, y + 0.7, 3.4, 0.5, title, font_size=18, bold=True, color=WHITE)
    
    # 内容
    add_bullet_list(slide, x + 0.2, y + 1.3, 3.4, 1.5, items, font_size=14, color=WHITE)

add_stage_card(slide, 0.5, 2.5, 1, "基础句型", [
    "主谓宾结构",
    "主系表结构",
    "疑问句构造",
    "时态入门"
], PRIMARY)

add_stage_card(slide, 4.7, 2.5, 2, "核心语法", [
    "8大核心时态",
    "冠词用法",
    "介词搭配",
    "代词指代"
], ACCENT)

add_stage_card(slide, 8.9, 2.5, 3, "进阶语法", [
    "从句结构",
    "非谓语动词",
    "虚拟语气",
    "复杂句式"
], RGBColor(0xB8, 0x50, 0x42))

add_text_box(slide, 0.5, 5.8, 12.3, 1, "学习形式", font_size=20, bold=True, color=PRIMARY)
add_bullet_list(slide, 0.5, 6.3, 12.3, 1.2, [
    "短句例句 + 图文讲解 → 填空练习 → 造句练习 → AI批改句子语法正误",
    "场景化拆解语法：把语法融入日常对话场景，学完立刻能用在口语里",
    "实时语法纠错：在用户写作或口语练习中，调用语法检查API纠正错误"
], font_size=16)

# 8. 学习路径
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_section_header(slide, "学习路径设计")

add_text_box(slide, 0.5, 1.8, 12.3, 1, "从零基础到自由沟通的完整路径", font_size=24, bold=True, color=PRIMARY)

# 创建学习路径时间轴
def add_path_stage(slide, x, y, stage, duration, title, items, color):
    # 阶段圆圈
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(1.5), Inches(1.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    
    # 阶段文字
    add_text_box(slide, x + 0.2, y + 0.3, 1.1, 0.4, f"阶段{stage}", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, y + 0.7, 1.1, 0.4, duration, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # 标题和内容
    add_text_box(slide, x - 0.5, y + 1.6, 2.5, 0.5, title, font_size=14, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)
    add_bullet_list(slide, x - 0.5, y + 2.1, 2.5, 1.5, items, font_size=12)

# 添加四个阶段
path_stages = [
    (1, "0-3个月", "音标发音筑基", ["搞定所有音标", "会拼读任意单词", "标准发音打底"], PRIMARY),
    (2, "3-6个月", "词汇+简单句型", ["背诵1500高频词", "掌握基础句式", "能说简单短句问答"], ACCENT),
    (3, "6-12个月", "系统语法+短句对话", ["吃透口语必备语法", "连贯多句对话", "应对所有日常场景"], RGBColor(0xB8, 0x50, 0x42)),
    (4, "持续", "自由沉浸式口语", ["AI无限话题闲聊", "角色扮演", "日常聊天、出行、职场"], RGBColor(0x2C, 0x5F, 0x2D))
]

for i, (stage, duration, title, items, color) in enumerate(path_stages):
    x = 1 + i * 3
    add_path_stage(slide, x, 2.5, stage, duration, title, items, color)
    # 连接箭头
    if i < len(path_stages) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x + 1.6), Inches(3.1), Inches(1.2), Inches(0.1)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY
        arrow.line.fill.background()

add_text_box(slide, 0.5, 5.5, 12.3, 1, "每日学习建议", font_size=20, bold=True, color=PRIMARY)
add_bullet_list(slide, 0.5, 6, 12.3, 1.5, [
    "发音跟读：15分钟",
    "单词记忆：20分钟（新词+复习）",
    "语法学习：15分钟",
    "综合练习：10分钟"
], font_size=16)

# 9. 技术架构
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_section_header(slide, "技术架构设计")

add_text_box(slide, 0.5, 1.8, 12.3, 1, "系统架构图", font_size=24, bold=True, color=PRIMARY)

# 创建架构图
def add_arch_box(slide, x, y, title, color, width=2.5, height=1):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    
    add_text_box(slide, x + 0.1, y + 0.2, width - 0.2, 0.6, title, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# 架构层级
add_arch_box(slide, 5, 2.5, "用户界面 (小程序)", PRIMARY, 3.3, 0.8)
add_arch_box(slide, 5, 3.8, "业务逻辑层 (云函数)", ACCENT, 3.3, 0.8)
add_arch_box(slide, 5, 5.1, "数据层 (云数据库)", RGBColor(0xB8, 0x50, 0x42), 3.3, 0.8)
add_arch_box(slide, 5, 6.4, "外部服务 (语音识别API、翻译API等)", RGBColor(0x2C, 0x5F, 0x2D), 3.3, 0.8)

# 连接箭头
for y in [3.3, 4.6, 5.9]:
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(y), Inches(0.1), Inches(0.5)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GRAY
    arrow.line.fill.background()

add_text_box(slide, 9, 2.5, 4, 1, "关键技术点", font_size=20, bold=True, color=PRIMARY)
add_bullet_list(slide, 9, 3, 4, 4, [
    "语音评测：微软Azure发音评测SDK",
    "AI对话：GPT-3.5/4 API或国产大模型",
    "数据存储：MySQL + Redis + OSS",
    "性能优化：CDN加速、离线缓存",
    "跨平台：UniApp/Taro一套代码",
    "推送服务：微信订阅消息"
], font_size=14)

# 10. 开发计划
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_section_header(slide, "开发计划")

add_text_box(slide, 0.5, 1.8, 12.3, 1, "三阶段开发计划", font_size=24, bold=True, color=PRIMARY)

# 创建三个阶段卡片
def add_dev_card(slide, x, y, phase, duration, title, items, color):
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(4.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = color
    card.line.width = Pt(2)
    
    # 阶段标签
    add_text_box(slide, x + 0.2, y + 0.2, 1.5, 0.5, f"第{phase}阶段", font_size=16, bold=True, color=color)
    
    # 时长
    add_text_box(slide, x + 1.8, y + 0.2, 2, 0.5, duration, font_size=14, color=GRAY)
    
    # 标题
    add_text_box(slide, x + 0.2, y + 0.8, 3.4, 0.5, title, font_size=18, bold=True, color=PRIMARY)
    
    # 内容
    add_bullet_list(slide, x + 0.2, y + 1.3, 3.4, 3, items, font_size=14)

add_dev_card(slide, 0.5, 2.5, "一", "1.5个月", "MVP上线", [
    "基础框架搭建（UniApp + 后端服务）",
    "音标跟读基础功能",
    "基础单词背诵（500词）",
    "入门语法课程（阶段1）",
    "基础AI短句对话",
    "打卡分享功能"
], PRIMARY)

add_dev_card(slide, 4.7, 2.5, "二", "上线后2个月", "功能完善", [
    "完整分级词库（1500词、3000词）",
    "全套语法体系（阶段2、3）",
    "多场景AI自由对话",
    "错题本、生词本",
    "会员付费体系"
], ACCENT)

add_dev_card(slide, 8.9, 2.5, "三", "用户破万后", "生态扩展", [
    "基于UniApp打包双端App",
    "打通小程序与App数据",
    "社区功能（学习圈、问答）",
    "多渠道运营",
    "持续迭代优化"
], RGBColor(0xB8, 0x50, 0x42))

# 11. 风险与对策
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, LIGHT_BG)
add_section_header(slide, "风险与对策")

add_text_box(slide, 0.5, 1.8, 12.3, 1, "潜在挑战与解决方案", font_size=24, bold=True, color=PRIMARY)

# 创建风险对策表格
def add_risk_table(slide, x, y):
    # 表格背景
    table_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(12.3), Inches(4)
    )
    table_bg.fill.solid()
    table_bg.fill.fore_color.rgb = WHITE
    table_bg.line.color.rgb = SECONDARY
    table_bg.line.width = Pt(1)
    
    # 表头
    headers = ["风险", "描述", "对策"]
    for i, header in enumerate(headers):
        add_text_box(slide, x + 0.2 + i * 4.1, y + 0.2, 3.9, 0.5, header, font_size=16, bold=True, color=PRIMARY)
    
    # 内容
    risks = [
        ("语音识别精度", "不同口音、环境噪音影响识别", "提供标准发音参考，允许重录；后期引入降噪算法"),
        ("内容版权", "单词库、例句可能涉及版权问题", "使用开源词典（ECdict），自编例句"),
        ("用户留存", "学习类应用容易流失", "游戏化设计（积分、勋章）、社交激励（组队PK）、定期推送提醒")
    ]
    
    for i, (risk, desc, solution) in enumerate(risks):
        y_pos = y + 1 + i * 1
        add_text_box(slide, x + 0.2, y_pos, 3.9, 0.8, risk, font_size=14, bold=True, color=PRIMARY)
        add_text_box(slide, x + 4.3, y_pos, 3.9, 0.8, desc, font_size=14)
        add_text_box(slide, x + 8.4, y_pos, 3.9, 0.8, solution, font_size=14, color=ACCENT)

add_risk_table(slide, 0.5, 2.5)

# 12. 总结与展望
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, PRIMARY)
add_text_box(slide, 1, 2, 11.3, 2, "总结与展望", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1, 4, 11.3, 1.5, "EnglishMaster - 让英语学习更简单、更高效、更有趣", font_size=24, color=SECONDARY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 2, 5.5, 9.3, 2, "核心优势：\n• AI纠音：音素级发音评估，纠正中式发音\n• 科学记忆：艾宾浩斯记忆曲线，高效背单词\n• 场景化语法：融入日常对话，学完就能用\n• AI对话陪练：模拟真实场景，锻炼口语反应", font_size=18, color=WHITE, alignment=PP_ALIGN.LEFT)

add_text_box(slide, 1, 7, 11.3, 0.5, "期待与您共同打造这款英语学习产品！", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# 保存演示文稿
prs.save("D:\\Desktop\\learnEngish\\英语学习小程序方案.pptx")
print("PPT已成功创建：英语学习小程序方案.pptx")